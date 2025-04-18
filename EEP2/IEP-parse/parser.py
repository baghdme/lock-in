from flask import Flask, request, jsonify, send_file
import os
import logging
from PIL import Image
import pytesseract
import pdf2image
import docx
from pdfminer.high_level import extract_text
from pptx import Presentation
import io
import json
from datetime import datetime
from sqlalchemy import create_engine, Column, String, DateTime, JSON, Text
from sqlalchemy.orm import declarative_base, sessionmaker, scoped_session
from sqlalchemy.sql import func
from dotenv import load_dotenv
import uuid

# Load environment variables
load_dotenv()

app = Flask(__name__)

# Configure logging
logging.basicConfig(level=logging.INFO)
logger = logging.getLogger(__name__)

# Initialize SQLAlchemy
Base = declarative_base()

class ParsedDocument(Base):
    """Schema for parsed documents"""
    __tablename__ = 'parsed_documents'
    
    id = Column(String, primary_key=True)
    original_filename = Column(String, nullable=False)
    file_type = Column(String, nullable=False)
    text_content = Column(Text, nullable=False)
    doc_metadata = Column(JSON, default={})
    created_at = Column(DateTime, default=func.now())

# Initialize database
db_path = os.path.join(os.path.dirname(__file__), 'iep_parse.db')
engine = create_engine(
    f'sqlite:///{db_path}',
    echo=True
)
Base.metadata.create_all(engine)
Session = scoped_session(sessionmaker(bind=engine))

# Configure Tesseract path
tesseract_path = os.getenv('TESSERACT_PATH')
if tesseract_path:
    pytesseract.pytesseract.tesseract_cmd = tesseract_path

def save_parsed_document(filename: str, file_type: str, text_content: str, metadata: dict = None):
    """Save parsed document to database"""
    session = Session()
    try:
        doc = ParsedDocument(
            id=str(uuid.uuid4()),
            original_filename=filename,
            file_type=file_type,
            text_content=text_content,
            doc_metadata=metadata or {}
        )
        session.add(doc)
        session.commit()
        # Refresh the document to get the created_at timestamp
        session.refresh(doc)
        # Create a dictionary with the document data before closing the session
        doc_data = {
            'id': doc.id,
            'original_filename': doc.original_filename,
            'file_type': doc.file_type,
            'text_content': doc.text_content,
            'doc_metadata': doc.doc_metadata,
            'created_at': doc.created_at
        }
        return doc_data
    except Exception as e:
        session.rollback()
        raise e
    finally:
        session.close()

def get_parsed_document(doc_id: str):
    """Get parsed document from database"""
    session = Session()
    try:
        return session.query(ParsedDocument).filter_by(id=doc_id).first()
    finally:
        session.close()

def extract_text_from_pdf(file):
    """Extract text from PDF file"""
    try:
        # Save the file temporarily
        temp_path = os.path.join(os.path.dirname(__file__), 'temp.pdf')
        file.save(temp_path)
        
        # Extract text using pdfminer
        text = extract_text(temp_path)
        
        # Clean up
        os.remove(temp_path)
        
        return text
    except Exception as e:
        logger.error(f"Error extracting text from PDF: {str(e)}")
        raise

def extract_text_from_docx(file):
    """Extract text from DOCX file"""
    try:
        doc = docx.Document(file)
        text = '\n'.join([paragraph.text for paragraph in doc.paragraphs])
        return text
    except Exception as e:
        logger.error(f"Error extracting text from DOCX: {str(e)}")
        raise

def extract_text_from_pptx(file):
    """Extract text from PPTX file"""
    try:
        prs = Presentation(file)
        text_runs = []
        
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text_runs.append(shape.text)
        
        return '\n'.join(text_runs)
    except Exception as e:
        logger.error(f"Error extracting text from PPTX: {str(e)}")
        raise

@app.route('/parse', methods=['POST'])
def parse_document():
    """Parse uploaded document and extract text"""
    if 'file' not in request.files:
        return jsonify({'error': 'No file part'}), 400
    
    file = request.files['file']
    if file.filename == '':
        return jsonify({'error': 'No selected file'}), 400
    
    try:
        filename = file.filename
        file_type = filename.rsplit('.', 1)[1].lower()
        
        # Extract text based on file type
        if file_type == 'pdf':
            text = extract_text_from_pdf(file)
        elif file_type in ['doc', 'docx']:
            text = extract_text_from_docx(file)
        elif file_type == 'pptx':
            text = extract_text_from_pptx(file)
        else:
            return jsonify({'error': 'Unsupported file type'}), 400
        
        # Save parsed document
        doc_data = save_parsed_document(
            filename=filename,
            file_type=file_type,
            text_content=text,
            metadata={'parse_time': datetime.utcnow().isoformat()}
        )
        
        return jsonify({
            'id': doc_data['id'],
            'text': doc_data['text_content'],
            'metadata': doc_data['doc_metadata']
        })
        
    except Exception as e:
        logger.error(f"Error parsing document: {str(e)}")
        return jsonify({'error': 'Error processing file'}), 500

@app.route('/documents/<doc_id>', methods=['GET'])
def get_document(doc_id):
    """Get parsed document by ID"""
    doc = get_parsed_document(doc_id)
    if not doc:
        return jsonify({'error': 'Document not found'}), 404
    
    return jsonify({
        'id': doc.id,
        'filename': doc.original_filename,
        'file_type': doc.file_type,
        'text': doc.text_content,
        'metadata': doc.doc_metadata,
        'created_at': doc.created_at.isoformat()
    })

@app.route('/set-tesseract-path', methods=['POST'])
def set_tesseract_path():
    """Set Tesseract executable path"""
    data = request.json
    if not data or 'path' not in data:
        return jsonify({'error': 'No path provided'}), 400
    
    try:
        pytesseract.pytesseract.tesseract_cmd = data['path']
        return jsonify({'message': 'Tesseract path updated successfully'})
    except Exception as e:
        logger.error(f"Error setting Tesseract path: {str(e)}")
        return jsonify({'error': 'Error setting Tesseract path'}), 500

@app.route('/health', methods=['GET'])
def health_check():
    """Health check endpoint."""
    try:
        # Check database
        session = Session()
        session.execute('SELECT 1')
        session.close()
        
        # Test Tesseract
        test_image = Image.new('RGB', (100, 30), color='white')
        pytesseract.image_to_string(test_image)
        
        return jsonify({
            'status': 'healthy',
            'services': {
                'database': 'healthy',
                'tesseract': 'healthy'
            }
        })
    except Exception as e:
        logger.error(f"Health check failed: {str(e)}")
        return jsonify({
            'status': 'unhealthy',
            'error': str(e)
        }), 500

if __name__ == '__main__':
    app.run(host='0.0.0.0', port=5003, debug=True)
